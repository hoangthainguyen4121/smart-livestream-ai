"""Extract text from thesis course materials for review."""
from pathlib import Path

import docx
import fitz
from pptx import Presentation

BASE = Path(r"D:\Thac Si NTTU\Chuyen de cntt\Document")
OUT_DIR = Path(__file__).resolve().parents[1] / "docs" / "_extracted_course_materials"
OUT_DIR.mkdir(parents=True, exist_ok=True)


def extract_docx(path: Path) -> str:
    document = docx.Document(path)
    parts: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    presentation = Presentation(path)
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"--- Slide {index} ---\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def extract_pdf(path: Path) -> str:
    document = fitz.open(path)
    parts: list[str] = []
    for index, page in enumerate(document, 1):
        text = page.get_text().strip()
        if text:
            parts.append(f"--- Page {index} ---\n{text}")
    return "\n\n".join(parts)


def main() -> None:
    for path in sorted(BASE.iterdir()):
        suffix = path.suffix.lower()
        if suffix == ".docx":
            text = extract_docx(path)
        elif suffix == ".pptx":
            text = extract_pptx(path)
        elif suffix == ".pdf":
            text = extract_pdf(path)
        else:
            continue

        out_path = OUT_DIR / f"{path.stem}.txt"
        out_path.write_text(text, encoding="utf-8")


if __name__ == "__main__":
    main()
