# -*- coding: utf-8 -*-
"""Generate Presentation 1 PPTX from presentation1_slide_data."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from presentation1_slide_data import SLIDES, TITLE_THESIS

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "thesis-figures"
OUT = ROOT / "Document" / "Thuyet trinh 1 - Smart Livestream.pptx"

FONT = "Times New Roman"


def _font(run, size=18, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold


def _bullets(tf, lines: list[str], max_n=6):
    tf.clear()
    for i, line in enumerate(lines[:max_n]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            _font(run)


def _table(slide, rows, left=0.4, top=3.6, width=5.0, height=2.2):
    if not rows:
        return
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = shape.table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _font(run, 11, bold=(r == 0))


def _image(slide, filename, left=5.3, top=1.1, width=4.2):
    path = FIG / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def _notes(slide, text: str, transition: str):
    body = text
    if transition and transition != "—":
        body += f"\n\n[Chuyển tiếp] {transition}"
    tf = slide.notes_slide.notes_text_frame
    tf.text = body
    for p in tf.paragraphs:
        for run in p.runs:
            _font(run, 12)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        if spec["id"] == "S1":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = spec["title"]
            sub = slide.placeholders[1]
            sub.text = f"{TITLE_THESIS}\n[Thuyết trình 1 · Họ tên · GVHD · Ngày]"
            for p in sub.text_frame.paragraphs:
                for run in p.runs:
                    _font(run, 16)
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                _font(run, 26, True)
            _notes(slide, spec["notes"], spec["transition"])
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = spec["title"]
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            _font(run, 26, True)

        has_img = bool(spec.get("image"))
        has_tbl = bool(spec.get("table"))
        body = slide.placeholders[1]
        _bullets(body.text_frame, spec["bullets"])

        if has_img:
            _image(slide, spec["image"], left=5.0 if has_tbl else 5.3, width=4.0 if has_tbl else 4.3)
        if has_tbl:
            _table(slide, spec["table"], top=3.5 if has_img else 1.6, left=0.4, width=4.8 if has_img else 8.5)

        _notes(slide, spec["notes"], spec["transition"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT} ({len(SLIDES)} slides)")
    return OUT


if __name__ == "__main__":
    build()
